"""TODO:"""

from abc import ABC, abstractmethod
from os import listdir
from os.path import isfile, join, isdir
import PyPDF2
import docx
from pptx import Presentation



class Document(ABC):

    """
    TODO:
    """

    def __init__(self, file_obj):

        self.file_obj = file_obj

    def get_files(self, folder_path:str):

        """
        TODO:
        """
        file_list = []

        for folder_object in listdir(folder_path):

            folder_object_path = join(folder_path, folder_object)

            if isfile(folder_object_path):

                file_list.append(folder_object_path)

            elif isdir(folder_object):

                file_list = file_list + self.get_files(folder_path=folder_object)

        return file_list


    @classmethod
    @abstractmethod
    def read(cls) -> str:

        """
        TODO:
        """


class Txt(Document):

    """
    TODO:
    """

    def read(self) -> str:
        """
        TODO:
        """
        with open(self.file_obj, 'r', encoding="utf-8") as file_object:

            file_content = file_object.read()

        return file_content


class PDF(Document):
    """
    TODO:
    """

    def read(self) -> str:

        """
        TODO:
        """
        reader = PyPDF2.PdfReader(self.file_obj)
        file_content = ""
        pages_total = len(reader.pages)

        for page_i in range(0, pages_total):

            file_content = file_content + " " + reader.pages[page_i].extract_text()

        return file_content


class Word(Document):
    """
    TODO:
    """

    def read(self) -> str:
        """
        TODO:
        """
        reader = docx.Document(self.file_obj)
        file_content = []

        for paragraph_i in reader.paragraphs:

            file_content.append(paragraph_i.text)

        return "\n".join(file_content)


class PowerPoint(Document):
    """
    TODO:
    """
    def read(self) -> str:
        """
        TODO:
        """
        reader = Presentation(self.file_obj)
        file_content = []

        for slide_i in reader.slides:

            for shape_i in slide_i.shapes:

                file_content.append(shape_i.text)

            return "\n".join(file_content)


class DocumentFactory:

    """
    TODO:
    """


    def get_doc_obj(self, file_obj):
        """
        TODO:
        """
        if file_obj.endswith(".pdf"):
            return PDF(file_obj=file_obj)

        elif file_obj.endswith(".txt"):
            return Txt(file_obj=file_obj)

        elif file_obj.endswith(".docx"):
            return Word(file_obj=file_obj)

        elif file_obj.endswith(".pptx"):
            return PowerPoint(file_obj=file_obj)

        else:
            return None
